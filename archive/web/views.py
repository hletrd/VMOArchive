from django.shortcuts import render, redirect
from django.http import HttpResponse, FileResponse, HttpResponseRedirect, JsonResponse
from django.core.paginator import Paginator

from django.contrib.auth.decorators import login_required, permission_required, user_passes_test
from django.contrib.auth import authenticate, login, logout
from django.core.files import File
from django.db.models import Q
from django.conf import settings
from django.template.loader import render_to_string

from . import forms
from . import models

import random, string
import datetime
import os
import re
import requests
import subprocess
import base64
from PIL import Image, ImageDraw, ImageFont
import sqlite3
import io
import math

import PyPDF2
from pptx import Presentation

import threading, ssl, smtplib
from email.mime.text import MIMEText
from email.mime.base import MIMEBase
from email.mime.multipart import MIMEMultipart
from email.header import Header
from email import encoders

import pymysql.cursors

# Create your views here.
def index(request):
  return show_list(request)

def user_is_staff(user):
  return user.is_staff

def log_smtp(log):
  with open('result.txt', 'a') as f:
    f.write(str(log))
    f.write("\n")

class SendMail(threading.Thread):
  def __init__(self, title, content, to, cc=[], files=[]):
    threading.Thread.__init__(self)
    #msg = MIMEText(content.encode('utf-8'), 'html', 'utf-8')
    msg = MIMEMultipart()
    msg['Subject'] = Header(title, 'utf-8')
    msg['From'] = settings.MAIL_ID + '@' + settings.MAIL_DOMAIN
    msg['To'] = to
    msg['CC'] = ','.join(cc)
    msg.attach(MIMEText(content, 'html'))

    toaddrs = ','.join([to] + cc)
    #self.msg = msg

    for i in files:
      file = MIMEBase('application', 'octet-stream')
      file.set_payload(i['data'].encode())
      encoders.encode_base64(file)
      file.add_header('Content-Disposition', 'attachment; filename={}'.format(i['name']))
      msg.attach(file)

    try:
      context = ssl.SSLContext(ssl.PROTOCOL_TLS)
      result = connection = smtplib.SMTP(settings.MAIL_SMTP, settings.MAIL_PORT)
      result = connection.set_debuglevel(2)
      result = connection.ehlo()
      #result = connection.starttls(context=context)
      result = connection.starttls()
      result = connection.ehlo()
      result = connection.login(settings.MAIL_ID + '@' + settings.MAIL_DOMAIN, settings.MAIL_PW)
      result = connection.sendmail(msg['From'], toaddrs, msg.as_string())
      result = connection.quit()
    except Exception as e:
      log_smtp(e)

  def run(self):
    pass


def reminder(request):
  if request.GET.get('key', '') == settings.KEY:
    members = models.VMOMember.objects.all()
    recipients = []
    for i in members:
      if re.match(r'[^@]+@[^.]+\..*', i.email_reminder):
        recipients.append(i.email_reminder)

    title = 'Seminar Reminder (' + datetime.datetime.now().strftime("%-m월 %-d일 기준") + ')'

    dates = ['일', '월', '화', '수', '목', '금', '토']

    seminar = models.VMOSeminar.objects.filter(date__gte=datetime.datetime.now())[:5]
    statussharing = list(models.VMOStatusSharing.objects.filter(date__gte=datetime.datetime.now())[:10])

    seminar = list(map(lambda x: {'id': x.id, 'member': x.member, 'date': x.date.strftime("%-m월 %-d일 (") + dates[int(x.date.strftime("%w"))] + ")", 'date_raw': x.date}, seminar))
    statussharing = list(map(lambda x: {'id': x.id, 'member': x.member, 'date': x.date.strftime("%-m월 %-d일 (") + dates[int(x.date.strftime("%w"))] + ")", 'date_raw': x.date}, statussharing))

    for i in statussharing:
      cal = models.VMOCalendar.objects.filter(targettype=0, targetid=i['id'])
      if not cal.exists():
        starttime = datetime.datetime.strptime(settings.TIME_STATUS, "%H:%M:%S").timetz()
        endtime = datetime.datetime.strptime(settings.TIME_STATUS, "%H:%M:%S") + datetime.timedelta(hours=1)
        calendar = models.VMOCalendar(calendarid=create_random_path(), targettype=0, targetid=i['id'], member=i['member'], date=i['date_raw'], starttime=starttime, endtime=endtime)
        calendar.save()
        i['cal'] = calendar.calendarid
      else:
        i['cal'] = cal.get().calendarid

    for i in seminar:
      cal = models.VMOCalendar.objects.filter(targettype=1, targetid=i['id'])
      if not cal.exists():
        starttime = datetime.datetime.strptime(settings.TIME_SEMINAR, "%H:%M:%S").timetz()
        endtime = datetime.datetime.strptime(settings.TIME_SEMINAR, "%H:%M:%S") + datetime.timedelta(hours=1)
        calendar = models.VMOCalendar(calendarid=create_random_path(), targettype=1, targetid=i['id'], member=i['member'], date=i['date_raw'], starttime=starttime, endtime=endtime)
        calendar.save()
        i['cal'] = calendar.calendarid
      else:
        i['cal'] = cal.get().calendarid

    files = []
    for i in statussharing:
      cal = {}
      cal['name'] = 'calendar_{}_{}.ics'.format(0, i['id'])
      dtstamp = datetime.datetime.now().strftime("%Y%m%dT%H%M%S")
      dtstart = datetime.datetime.combine(i['date_raw'], datetime.datetime.strptime(settings.TIME_STATUS, "%H:%M:%S").timetz())
      dtstart = dtstart.strftime("%Y%m%dT%H%M%S")
      dtend = datetime.datetime.combine(i['date_raw'], (datetime.datetime.strptime(settings.TIME_STATUS, "%H:%M:%S") + datetime.timedelta(hours=1)).timetz())
      dtend = dtend.strftime("%Y%m%dT%H%M%S")
      caltitle = i['member'].name + '의 '
      caltitle += 'Status sharing'
      description = caltitle
      cal['data'] = render_to_string('calendar.ics', {'dtstamp': dtstamp, 'dtstart': dtstart, 'dtend': dtend, 'title': caltitle, 'description': description, 'calid': create_random_path()})
      files.append(cal)

    for i in seminar:
      cal = {}
      cal['name'] = 'calendar_{}_{}.ics'.format(1, i['id'])
      dtstamp = datetime.datetime.now().strftime("%Y%m%dT%H%M%S")
      dtstart = datetime.datetime.combine(i['date_raw'], datetime.datetime.strptime(settings.TIME_SEMINAR, "%H:%M:%S").timetz())
      dtstart = dtstart.strftime("%Y%m%dT%H%M%S")
      dtend = datetime.datetime.combine(i['date_raw'], (datetime.datetime.strptime(settings.TIME_SEMINAR, "%H:%M:%S") + datetime.timedelta(hours=1)).timetz())
      dtend = dtend.strftime("%Y%m%dT%H%M%S")
      caltitle = i['member'].name + '의 '
      caltitle += 'Seminar'
      description = caltitle
      cal['data'] = render_to_string('calendar.ics', {'dtstamp': dtstamp, 'dtstart': dtstart, 'dtend': dtend, 'title': caltitle, 'description': description, 'calid': create_random_path()})
      files.append(cal)

    status_origin = datetime.datetime.strptime(settings.STATUS_ORIGIN, "%Y-%m-%d")
    now = datetime.datetime.now()
    delta = now - status_origin
    due = status_origin + datetime.timedelta(days=math.ceil(delta.days/14)*14)
    due = due.strftime("%-m/%-d (") + dates[int(due.strftime("%w"))] + ')'

    content = render_to_string('reminder.html', {'seminar': seminar, 'statussharing': statussharing, 'due': due})

    result = 'Sent mail: <br/><br/>' + content
    result += '<br/><br/>To: '
    
    to = settings.EMAIL_MASS
    cc = [settings.EMAIL_PROF]

    SM = SendMail(title, content, to, cc, files=files)
    result += to
    #for i in recipients:
    #  SM = SendMail(title, content, i)
      #SM.start()
    #  result += i + ', '
    return HttpResponse(result)
  else:
    return redirect('/')

def status_report(request):
  if request.GET.get('key', '') == settings.KEY:
    members = models.VMOMember.objects.all()
    recipients = []
    for i in members:
      if re.match(r'[^@]+@[^.]+\..*', i.email_reminder):
        recipients.append(i.email_reminder)

    dates = ['일', '월', '화', '수', '목', '금', '토']

    status_origin = datetime.datetime.strptime("2022-03-18", "%Y-%m-%d")
    now = datetime.datetime.now()
    delta = now - status_origin
    due = status_origin + datetime.timedelta(days=math.ceil(delta.days/14)*14)

    if not (datetime.timedelta(days=0) <= (due - datetime.datetime.now()) < datetime.timedelta(days=1)):
      #pass
      return HttpResponse('no')

    due = due.strftime("%-m/%-d (") + dates[int(due.strftime("%w"))] + ')'
    title = 'Status report due ' + due

    seminar = models.VMOSeminar.objects.filter(date__gte=datetime.datetime.now())[:5]
    statussharing = list(models.VMOStatusSharing.objects.filter(date__gte=datetime.datetime.now())[:10])

    seminar = list(map(lambda x: {'member': x.member, 'date': x.date.strftime("%-m월 %-d일 (") + dates[int(x.date.strftime("%w"))] + ")"}, seminar))
    statussharing = list(map(lambda x: {'member': x.member, 'date': x.date.strftime("%-m월 %-d일 (") + dates[int(x.date.strftime("%w"))] + ")"}, statussharing))

    content = render_to_string('status_report.html', {'seminar': seminar, 'statussharing': statussharing, 'due': due})

    result = 'Sent mail: <br/><br/>' + content
    result += '<br/><br/>To: '
    
    to = settings.EMAIL_MASS
    cc = [settings.EMAIL_PROF]

    SM = SendMail(title, content, to, cc)
    result += to
    #for i in recipients:
    #  SM = SendMail(title, content, i)
      #SM.start()
    #  result += i + ', '
    return HttpResponse(result)
  else:
    return redirect('/')

@login_required
def calendar(request, calendarid):
  cal = models.VMOCalendar.objects.get(calendarid=calendarid)
  if cal != None:
    dtstamp = datetime.datetime.now().strftime("%Y%m%dT%H%M%S")
    dtstart = datetime.datetime.combine(cal.date, cal.starttime)
    dtstart = dtstart.strftime("%Y%m%dT%H%M%S")
    dtend = datetime.datetime.combine(cal.date, cal.endtime)
    dtend = dtend.strftime("%Y%m%dT%H%M%S")
    title = cal.member.name + '의 '
    if cal.targettype == 1:
      title += 'Seminar'
    if cal.targettype == 0:
      title += 'Status sharing'
    
    description = title
    result = render_to_string('calendar.ics', {'dtstamp': dtstamp, 'dtstart': dtstart, 'dtend': dtend, 'title': title, 'description': description, 'calid': calendarid})
    return HttpResponse(result, headers={'Content-Type': 'text/calendar', 'Content-Disposition': 'inline; filename="calendar.ics"'})
  else:
    return redirect('/')

def calendar_add(request, calendarid):
  cal = models.VMOCalendar.objects.get(calendarid=calendarid)
  if cal != None:
    HttpResponseRedirect.allowed_schemes.append('webcal')
    return redirect('webcal://vmo.snu.ac.kr:8888/calendar/{}'.format(calendarid))
  else:
    return redirect('/')

@login_required
def mass(request):
	if request.method == 'POST':
		form = forms.mail_form(request.POST)
		if form.is_valid():
			if form.cleaned_data.get('recipient') == 'mass@altair.snu.ac.kr':
				recipients = [settings.EMAIL_PROF]
			elif form.cleaned_data.get('recipient') == 'masss@altair.snu.ac.kr':
				recipients = []
			else:
				return render(request, 'mail_form.html', {'error': '잘못된 입력값이 있습니다.', 'range': form.cleaned_data.get('range_', 2), 'subject': form.cleaned_data.get('subject', ''), 'content': form.cleaned_data.get('content', '')})
			members = models.VMOMember.objects.all()
			for i in members:
				if re.match(r'[^@]+@[^.]+\..*', i.email_reminder):
					recipients.append(i.email_reminder)

			subject = form.cleaned_data.get('subject', '(No subject)')
			content = form.cleaned_data.get('content', '(No content)').replace("\n", "<br/>")

			for i in recipients:
				SM = SendMail(subject, content, i)
				SM.start()
			return render(request, 'jumbotron.html', {'output': 'Mail successfully sent to: {}'.format(', '.join(recipients))})
		else:
			return render(request, 'mail_form.html', {'error': '잘못된 입력값이 있습니다.', 'range': form.cleaned_data.get('range_', 2), 'subject': form.cleaned_data.get('subject', ''), 'content': form.cleaned_data.get('content', '')})
	else:
		return render(request, 'mail_form.html', {'range': 2})
def server_update(request):
  if request.GET.get('key', '') != settings.KEY:
    return redirect('/')

  conn = sqlite3.connect(settings.DB_SERVER)
  c = conn.cursor()

  c.execute('CREATE TABLE IF NOT EXISTS `status`(`id` INTEGER PRIMARY KEY AUTOINCREMENT, `servername` TEXT, `status` INT, `time` DATETIME);')
  c.execute('CREATE INDEX IF NOT EXISTS `servername` ON `status` (`servername`);')
  conn.commit()

  for i in settings.SERVERS:
    command = ['/usr/bin/fping', '-c1', '-t500', i+'.snu.ac.kr']
    status = subprocess.call(command) == 0
    c.execute("INSERT INTO `status`(`servername`, `status`, `time`) VALUES (?,?,datetime('now', 'localtime'));", (i, status))
  conn.commit()
  return HttpResponse('')

def create_graph(servername, width=240, height=60, hours=24):
  conn = sqlite3.connect(settings.DB_SERVER)
  c = conn.cursor()

  now = datetime.datetime.now()

  c.execute("SELECT * FROM `status` WHERE `servername`=? AND `time`>datetime('now', '-" + str(hours) + " hours');", (servername, ))
  data = c.fetchall()

  result = Image.new("RGBA", (hours*60, height), (255, 255, 255, 0))
  d = ImageDraw.Draw(result)

  color = {
    0: '#dc3545',
    1: '#20c997'
  }

  for i in data:
    date = datetime.datetime.strptime(i[3], '%Y-%m-%d %H:%M:%S')
    delta = now - date
    seconds = delta.total_seconds()
    mins = seconds / 60

    xy = [hours*60 - mins, 0, hours*60, height]
    d.rectangle(xy, fill=color[i[2]])

  offset = now.minute
  left = now + datetime.timedelta(hours=-hours)
  for i in range(hours+1):
    time = left + datetime.timedelta(hours=i)
    if time.hour%24 == 0:
      d.line([-offset+i*60, 0, -offset+i*60, height], '#555', width=2)
    elif time.hour%24 == 6 or time.hour%24 == 12 or time.hour%24 == 18:
      d.line([-offset+i*60, 0, -offset+i*60, height], '#888', width=2)
    else:
      d.line([-offset+i*60, 0, -offset+i*60, height], '#fff')

  font = ImageFont.truetype('../opensans.ttf', 17)

  result = result.resize([width, height], Image.BICUBIC)

  box_w = 120
  box_h = 120
  for i in range(hours+1):
    time = left + datetime.timedelta(hours=i)
    if time.hour % 6 != 0: continue
    temp = Image.new("RGBA", (box_w, box_h), (255, 255, 255, 0))
    drawtemp = ImageDraw.Draw(temp)
    drawtemp.text((box_w/2, box_h/2), '{:02d}:00'.format(time.hour), '#333', font, anchor='mm')
    temp = temp.rotate(90)
    result.paste(temp, (int(
      -offset/60/hours*width #minute offset
      +i/hours*width #hour positioning
      #+width/hours*1 #width compensation
      -box_w*0.5 #accounting bounding box size
      -1), int(-box_h*0.5+30)), temp)

  buffer = io.BytesIO()
  result.save(buffer, 'png')
  return buffer.getvalue()


from threading import Thread
def ping(index, result, address):
  command = ['/usr/bin/fping', '-c1', '-t100', address]
  result[index] = subprocess.call(command) == 0

@login_required
def server(request):
  width = 320
  height = 30
  hours = 48

  result = []
  result_ping = [False for i in range(len(settings.SERVERS))]
  threads = []
  for k, host in enumerate(settings.SERVERS):
    t = Thread(target=ping, args=(k, result_ping, host+'.snu.ac.kr'))
    t.start()
    threads.append(t)
    graph = create_graph(host, width*2, height*2, hours)
    graph = base64.b64encode(graph).decode('ascii')
    result.append({
      'name': host,
      'status': None,
      'graph': graph
    })
  for i in threads:
    i.join()
  for i in range(len(settings.SERVERS)):
    result[i]['status'] = result_ping[i]
  return render(request, 'server.html', {'result': result, 'controllable': settings.SERVERS_CONTROLLABLE, 'width': width, 'height': height})

@login_required
def server_action(request, action):
  servername = action.split('.')[0]
  do = action.split('.')[-1]
  if servername in settings.SERVERS_CONTROLLABLE:
    if do == 'short':
      r = requests.get(settings.SERVERS_CONTROLLABLE[servername] + '/run')
      result = r.text
    elif do == 'long':
      r = requests.get(settings.SERVERS_CONTROLLABLE[servername] + '/long')
      result = r.text
  return redirect('/')

def db_conn():
  conn = pymysql.connect(host='localhost', user=settings.DB_USER, password=settings.DB_PASS, database=settings.DB_NAME, charset='utf8mb4', cursorclass=pymysql.cursors.DictCursor)
  return conn.cursor()

def serverroom(request):
  c = db_conn()
  c.execute("SELECT * FROM `cache`;")
  now = c.fetchone()
  now = {
    'h': '{:.2f}'.format(now['h']),
    't': '{:.2f}'.format(now['t']),
  }

  return render(request, 'serverroom.html', {'now': now})

def serverroom_graph(request):
  c = db_conn()
  c.execute("SELECT UNIX_TIMESTAMP(`date`)*1000 AS `d`, ROUND(AVG(`h`),2) AS `h`, ROUND(AVG(`t`),2) AS `t` FROM `ht` WHERE `date` >= DATE_ADD(NOW(), INTERVAL -7 DAY) GROUP BY UNIX_TIMESTAMP(`date`) DIV 300 ORDER BY `id` DESC;")
  ht = c.fetchall()

  return JsonResponse(ht, safe=False)

def serverroom_get(request):
  c = db_conn()
  c.execute("SELECT * FROM `cache`;")
  now = c.fetchone()
  return JsonResponse({'date': now['date'].strftime("%Y-%m-%dT%H:%M:%S"), 'h': '{:.2f}'.format(now['h']), 't': '{:.2f}'.format(now['t'])})

@login_required
@user_passes_test(user_is_staff)
def upload(request):
	if request.method == 'POST':
		form = forms.upload_form(request.POST, request.FILES)
		if form.is_valid():
			post = models.VMODocument(title=form.cleaned_data['title'], content=form.cleaned_data['content'])
			post.datetime = form.cleaned_data['date']
			post.save()
			files = request.FILES.getlist('file')
			for i in files:
				path = create_random_path(10)
				file = models.VMOFile(path_raw=path, name=i.name, metadata='')
				i.name = path
				file.file = i
				file.save()
				create_index(file)
				post.files.add(file)
			post.save()
			return redirect('/')
		else:
			return render(request, 'upload.html', {'date': form.cleaned_data.get('date', datetime.datetime.now()).strftime('%Y-%m-%d'), 'title': form.cleaned_data.get('title', ''), 'content': form.cleaned_data.get('content', ''), 'error': '잘못된 입력값이 있습니다.', 'url': '/archive/upload'})
	else:
		return render(request, 'upload.html', {'date': datetime.datetime.now().strftime('%Y-%m-%d'), 'url': '/archive/upload'})

@login_required
@user_passes_test(user_is_staff)
def modify(request, postid):
	post = models.VMODocument.objects.get(id=postid)
	if post is not None:
		if request.method == 'POST':
			form = forms.upload_form(request.POST, request.FILES)
			if form.is_valid():
				post.title = form.cleaned_data['title']
				post.content = form.cleaned_data['content']
				post.datetime = form.cleaned_data['date']
				post.save()
				file_prev = request.POST.getlist('file_prev[]')
				for i in post.files.all():
					if not str(i.id) in file_prev:
						remove = post.files.get(id=i.id)
						print(remove.file.path)
						if os.path.isfile(remove.file.path):
							os.remove(remove.file.path)
						remove.delete()
					else:
						create_index(i)
				files = request.FILES.getlist('file')
				for i in files:
					path = create_random_path(10)
					file = models.VMOFile(path_raw=path, name=i.name, metadata='')
					i.name = path
					file.file = i
					file.save()
					create_index(file)
					post.files.add(file)
				post.save()
				return redirect('/archive')
			else:
				return render(request, 'upload.html', {'date': form.cleaned_data.get('date', datetime.datetime.now()).strftime('%Y-%m-%d'), 'title': form.cleaned_data.get('title', ''), 'content': form.cleaned_data.get('content', ''), 'error': '잘못된 입력값이 있습니다.', 'files': post.files, 'url': '/archive/modify/{}'.format(postid)})
		else:
			return render(request, 'upload.html', {'date': post.datetime.strftime('%Y-%m-%d'), 'title': post.title, 'content': post.content, 'files': post.files, 'url': '/archive/modify/{}'.format(postid)})
	else:
		return redirect('/archive')

@login_required
@user_passes_test(user_is_staff)
def delete(request, postid):
	post = models.VMODocument.objects.get(id=postid)
	if post is not None:
		if request.method == 'POST':
			for i in post.files.all():
				if os.path.isfile(i.file.path):
					os.remove(i.file.path)
				i.delete()
			post.delete()
			return redirect('/archive')
		else:
			return render(request, 'delete.html', {'name': post.title, 'url': '/archive/delete/{}'.format(postid)})
	else:
		return redirect('/archive') #404

@login_required
@user_passes_test(user_is_staff)
def init(request):
	import sqlite3
	import requests

	db = '../posts.db'
	conn = sqlite3.connect(db)

	models.VMODocument.objects.all().delete()
	models.VMOFile.objects.all().delete()

	c = conn.cursor()

	c.execute('SELECT * FROM posts ORDER BY `id` DESC;')
	posts = c.fetchall()

	for i in posts:
		content = ''
		if i[3] != 'None':
			content = i[3].rstrip()
		post = models.VMODocument(title=i[2], content=content)
		post.datetime = i[1]
		post.save()
		c.execute('SELECT * FROM files where postid=?', (i[0],))
		files = c.fetchall()
		for j in files:
			file = models.VMOFile(path_raw=j[2], name=j[1], metadata='')
			file.file = File(open('./files/' + j[2], 'rb'), name=j[2])
			file.save()
			post.files.add(file)
		post.save()
	return redirect('/archive')


@login_required
@user_passes_test(user_is_staff)
def index_file(request):
	files = models.VMOFile.objects.all()
	for i in files:
		create_index(i)
	return redirect('/')

def create_index(file):
	ext = file.name.split('.')[-1].lower()
	path = settings.MEDIA_ROOT + '/' + file.path_raw
	text = ''
	if ext == 'pptx':
		try:
			reader = Presentation(path)
			for slide in reader.slides:
				for shape in slide.shapes:
					if hasattr(shape, 'text'):
						text += shape.text + ' '
		except:
			print(file.path_raw)
			print(file.name)
			pass
	elif ext == 'pdf':
		reader = PyPDF2.PdfFileReader(path)
		for page in range(0, reader.getNumPages()):
			text += reader.getPage(page).extractText() + ' '
	file.metadata = text
	file.save()

@login_required
def show_list(request, page=1):
	search_text = request.GET.get('search', '')
	if search_text != '':
		search_inputs = search_text.split(' ')
		query = Q()
		for i in search_inputs:
			query |= Q(files__name__contains=i) | Q(title__contains=i) | Q(content__contains=i) | Q(files__metadata__contains=i)
		posts = models.VMODocument.objects.filter(query).order_by('-id')
	else:
		posts = models.VMODocument.objects.all().order_by('-id')

	page = int(page)
	paginator = Paginator(posts, 20)
	posts = paginator.get_page(page)

	num_pages = paginator.num_pages
	last_page = min(page+3, num_pages)+1
	pages = range(max(page-3, 1), last_page)

	for i in posts:
		i.content = i.content.replace('\n', '<br/>')

	return render(request, 'archive.html', {'posts': posts, 'pages': pages, 'num_pages': num_pages, 'last_page': last_page, 'search': search_text})

@login_required
def download(request, fileid):
	file = models.VMOFile.objects.get(path_raw=fileid)
	if file is not None:
		return FileResponse(file.file, filename=file.name)
	else:
		return redirect('/archive')

def create_random_path(length=10):
	path = ''.join(random.choice(string.ascii_lowercase + string.digits) for _ in range(length))
	return path

def handle_file(f):
	path = create_random_path(10)
	file = models.VMOFile.objects.create(path_raw=path, name=f.name, metadata='', file=f)
	#with open(settings.MEDIA_ROOT + '/' + path, 'wb+') as dest:
	#	for chunk in f.chunks():
	#		dest.write(chunk)
	return file

def login_view(request):
	if request.user.is_authenticated:
		return redirect('/')
	if request.method == 'POST':
		form = forms.login_form(request.POST)
		if form.is_valid():
			user = authenticate(request, username=form.cleaned_data['userid'], password=form.cleaned_data['password'])
			if user is not None:
				login(request, user)
				return redirect(request.GET.get('next', '/'))
			else:
				return render(request, 'login.html', {'error': '로그인에 실패했습니다.'})
		else:
			return render(request, 'login.html', {'error': '로그인에 실패했습니다.'})
	return render(request, 'login.html', {'next': request.GET.get('next')})

@login_required
def logout_view(request):
	logout(request)
	return redirect('/')

@login_required
def members(request):
	members = models.VMOMember.objects.all()
	return render(request, 'members.html', {'members': members})

@login_required
@user_passes_test(user_is_staff)
def members_add(request):
	if request.method == 'POST':
		form = forms.member_form(request.POST)
		if form.is_valid():
			member = models.VMOMember(name=form.cleaned_data['name'], birthday=form.cleaned_data['birthday'], phone=form.cleaned_data['phone'], email=form.cleaned_data['email'],email_reminder=form.cleaned_data['email_reminder'])
			member.save()
			return redirect('/members')
		else:
			return render(request, 'members_form.html', {'url': '/members/add', 'name': form.cleaned_data.get('name', ''), 'birthday': form.cleaned_data.get('birthday', ''), 'phone': form.cleaned_data.get('phone', ''), 'email': form.cleaned_data.get('email', ''), 'email_reminder': form.cleaned_data.get('email_reminder', ''), 'error': '잘못된 입력값이 있습니다.'})
	else:
		return render(request, 'members_form.html', {'url': '/members/add'})

@login_required
@user_passes_test(user_is_staff)
def members_modify(request, memberid):
	member = models.VMOMember.objects.get(id=memberid)
	if member is not None:
		if request.method == 'POST':
			form = forms.member_form(request.POST)
			if form.is_valid():
				member.name = form.cleaned_data['name']
				member.birthday = form.cleaned_data['birthday']
				member.phone = form.cleaned_data['phone']
				member.email = form.cleaned_data['email']
				member.email_reminder = form.cleaned_data['email_reminder']
				member.save()
				return redirect('/members')
			else:
				return render(request, 'members_form.html', {'url': '/members/modify/{}'.format(memberid), 'name': form.cleaned_data.get('name', ''), 'birthday': form.cleaned_data.get('birthday', ''), 'phone': form.cleaned_data.get('phone', ''), 'email': form.cleaned_data.get('email', ''), 'email_reminder': form.cleaned_data.get('email_reminder', ''), 'error': '잘못된 입력값이 있습니다.'})
		else:
			return render(request, 'members_form.html', {'url': '/members/modify/{}'.format(memberid), 'name': member.name, 'birthday': member.birthday.strftime("%m%d"), 'phone': member.phone, 'email': member.email, 'email_reminder': member.email_reminder})
	else:
		return redirect('/members')

@login_required
@user_passes_test(user_is_staff)
def members_delete(request, memberid):
	member = models.VMOMember.objects.get(id=memberid)
	if member is not None:
		if request.method == 'POST':
			member.delete()
			return redirect('/members')
		else:
			return render(request, 'delete.html', {'name': member.name, 'url': '/members/delete/{}'.format(memberid)})
	else:
		return redirect('/members')

#@login_required
def seminar(request):
	seminar = models.VMOSeminar.objects.all()
	return render(request, 'seminar.html', {'seminar': seminar, 'pagename': 'Seminar 일정', 'url': '/seminar'})

@login_required
@user_passes_test(user_is_staff)
def seminar_add(request):
	if request.method == 'POST':
		form = forms.seminar_form(request.POST)
		if form.is_valid():
			seminar = models.VMOSeminar(member=form.cleaned_data.get('member'), date=form.cleaned_data['date'])
			seminar.save()
			return redirect('/seminar')
		return render(request, 'seminar_form.html', {'form': form, 'url': '/seminar/add', 'member': form.cleaned_data.get('member', ''), 'date': form.cleaned_data.get('date', datetime.datetime.now()).strftime('%Y-%m-%d'), 'error': '잘못된 입력값이 있습니다.'})
	else:
		form = forms.seminar_form()
		return render(request, 'seminar_form.html', {'form': form, 'url': '/seminar/add'})

@login_required
@user_passes_test(user_is_staff)
def seminar_modify(request, seminarid):
	seminar = models.VMOSeminar.objects.get(id=seminarid)
	if seminar is not None:
		if request.method == 'POST':
			form = forms.seminar_form(request.POST, initial={'member': seminar.member})
			if form.is_valid():
				seminar.member = form.cleaned_data.get('member')
				seminar.date = form.cleaned_data['date']
				seminar.save()
				return redirect('/seminar')
			return render(request, 'seminar_form.html', {'form': form, 'url': '/seminar/modify/{}'.format(seminarid), 'member': form.cleaned_data.get('member', None), 'date': form.cleaned_data.get('date', datetime.datetime.now()).strftime('%Y-%m-%d'), 'error': '잘못된 입력값이 있습니다.'})
		else:
			form = forms.seminar_form(initial={'member': seminar.member})
			return render(request, 'seminar_form.html', {'form': form, 'url': '/seminar/modify/{}'.format(seminarid), 'member': seminar.member, 'date': seminar.date.strftime("%Y-%m-%d")})
	else:
		return redirect('/seminar')

@login_required
@user_passes_test(user_is_staff)
def seminar_delete(request, seminarid):
	seminar = models.VMOSeminar.objects.get(id=seminarid)
	if seminar is not None:
		if request.method == 'POST':
			seminar.delete()
			return redirect('/seminar')
		else:
			return render(request, 'delete.html', {'name': seminar.member.name + '의 ' + seminar.date.strftime("%m월 %d일 seminar"), 'url': '/seminar/delete/{}'.format(seminarid)})
	else:
		return redirect('/seminar')

@login_required
@user_passes_test(user_is_staff)
def seminar_change(request, day):
	seminar = models.VMOSeminar.objects.all()
	try:
		day = int(day)
	except:
		day = 0
	days = datetime.timedelta(days=day)
	for i in seminar:
		i.date += days
		i.save()
	return redirect('/seminar')

@login_required
@user_passes_test(user_is_staff)
def seminar_archive(request, seminarid):
	seminar = models.VMOSeminar.objects.get(id=seminarid)
	if seminar is not None:
		form = forms.upload_form()
		return render(request, 'upload.html', {'date': seminar.date.strftime("%Y-%m-%d"), 'title': 'Status Report {}'.format(seminar.date.strftime("%Y-%m-%d")), 'content': 'Status Report\n{}\n\n{}'.format(seminar.date.strftime("%Y-%m-%d"), seminar.member.name), 'url': '/archive/upload'})
	return redirect('/seminar')

#@login_required
def status_sharing(request):
	statussharing = models.VMOStatusSharing.objects.all()
	return render(request, 'seminar.html', {'seminar': statussharing, 'pagename': 'Status sharing 일정', 'url': '/status-sharing'})

@login_required
@user_passes_test(user_is_staff)
def status_sharing_add(request):
	if request.method == 'POST':
		form = forms.seminar_form(request.POST)
		if form.is_valid():
			status_sharing = models.VMOStatusSharing(member=form.cleaned_data.get('member'), date=form.cleaned_data['date'])
			status_sharing.save()
			return redirect('/status-sharing')
		return render(request, 'seminar_form.html', {'form': form, 'url': '/status-sharing/add', 'member': form.cleaned_data.get('member', ''), 'date': form.cleaned_data.get('date', datetime.datetime.now()).strftime('%Y-%m-%d'), 'error': '잘못된 입력값이 있습니다.'})
	else:
		form = forms.seminar_form()
		return render(request, 'seminar_form.html', {'form': form, 'url': '/status-sharing/add'})

@login_required
@user_passes_test(user_is_staff)
def status_sharing_modify(request, statussharingid):
	statussharing = models.VMOStatusSharing.objects.get(id=statussharingid)
	if statussharing is not None:
		if request.method == 'POST':
			form = forms.seminar_form(request.POST, initial={'member': statussharing.member})
			if form.is_valid():
				statussharing.member = form.cleaned_data.get('member')
				statussharing.date = form.cleaned_data['date']
				statussharing.save()
				return redirect('/status-sharing')
			return render(request, 'seminar_form.html', {'form': form, 'url': '/status-sharing/modify/{}'.format(statussharingid), 'member': form.cleaned_data.get('member', None), 'date': form.cleaned_data.get('date', datetime.datetime.now()).strftime('%Y-%m-%d'), 'error': '잘못된 입력값이 있습니다.'})
		else:
			form = forms.seminar_form(initial={'member': statussharing.member})
			return render(request, 'seminar_form.html', {'form': form, 'url': '/status-sharing/modify/{}'.format(statussharingid), 'member': statussharing.member, 'date': statussharing.date.strftime("%Y-%m-%d")})
	else:
		return redirect('/status-sharing')

@login_required
@user_passes_test(user_is_staff)
def status_sharing_delete(request, statussharingid):
	statussharing = models.VMOStatusSharing.objects.get(id=statussharingid)
	if statussharing is not None:
		if request.method == 'POST':
			statussharing.delete()
			return redirect('/status-sharing')
		else:
			return render(request, 'delete.html', {'name': statussharing.member.name + '의 ' + statussharing.date.strftime("%m월 %d일 status sharing"), 'url': '/status-sharing/delete/{}'.format(statussharingid)})
	else:
		return redirect('/status-sharing')

@login_required
@user_passes_test(user_is_staff)
def status_sharing_change(request, day):
	statussharing = models.VMOStatusSharing.objects.all()
	try:
		day = int(day)
	except:
		day = 0
	days = datetime.timedelta(days=day)
	for i in statussharing:
		i.date += days
		i.save()
	return redirect('/status-sharing')

@login_required
@user_passes_test(user_is_staff)
def status_sharing_archive(request, statussharingid):
	statussharing = models.VMOStatusSharing.objects.get(id=statussharingid)
	if statussharing is not None:
		form = forms.upload_form()
		return render(request, 'upload.html', {'date': statussharing.date.strftime("%Y-%m-%d"), 'title': '', 'content': 'VMOLAB Seminar\n{}\n\n'.format(statussharing.date.strftime("%Y-%m-%d")), 'url': '/archive/upload'})
	return redirect('/statussharing')